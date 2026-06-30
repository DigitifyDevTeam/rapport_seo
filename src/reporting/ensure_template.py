"""Ensure the PowerPoint deck matches current code (Backlinks, no Merci slide)."""

from __future__ import annotations

import logging
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from src.config import PROJECT_ROOT, TEMPLATE_PATH

logger = logging.getLogger(__name__)

# Keep in sync with scripts/build_template.py TEMPLATE_BUILD_VERSION.
EXPECTED_TEMPLATE_VERSION = "2026-06-v8-clarity-chart-spacing"
EXPECTED_SLIDE_COUNT = 15


def _version_file_for(pptx_path: Path) -> Path:
    return pptx_path.parent / ".template_build_version"


def _read_installed_version(pptx_path: Path) -> str | None:
    path = _version_file_for(pptx_path)
    if not path.is_file():
        return None
    try:
        return path.read_text(encoding="utf-8").strip() or None
    except OSError:
        return None


def _slide_text_blob(prs: Presentation) -> str:
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                parts.append(text)
    return "\n".join(parts).upper()


def template_structure_valid(pptx_path: Path) -> bool:
    """True when the deck has Backlinks slides and no closing Merci slide."""
    if not pptx_path.is_file():
        return False
    try:
        prs = Presentation(str(pptx_path))
    except Exception:  # noqa: BLE001
        return False
    if len(prs.slides) != EXPECTED_SLIDE_COUNT:
        return False
    blob = _slide_text_blob(prs)
    if "MERCI POUR VOTRE" in blob:
        return False
    if blob.count("BACKLINKS") < 2:
        return False
    if "SYNTH" not in blob and "SYNTHESE" not in blob.replace("È", "E"):
        return False
    return True


def template_is_current(pptx_path: Path) -> bool:
    if _read_installed_version(pptx_path) != EXPECTED_TEMPLATE_VERSION:
        return False
    return template_structure_valid(pptx_path)


def _rebuild_template(pptx_path: Path) -> None:
    script = PROJECT_ROOT / "scripts" / "build_template.py"
    if not script.is_file():
        raise FileNotFoundError(f"Missing template builder: {script}")
    cmd = [
        sys.executable,
        str(script),
        "--force",
        "--output",
        str(pptx_path),
    ]
    logger.info("Rebuilding report template → %s", pptx_path)
    subprocess.run(cmd, check=True, cwd=str(PROJECT_ROOT))


def ensure_report_template(pptx_path: Path | None = None) -> Path:
    """Rebuild ``pptx_path`` when missing, outdated, or still the old Merci deck."""
    target = Path(pptx_path or TEMPLATE_PATH).expanduser().resolve()
    if template_is_current(target):
        logger.debug("Report template OK: %s", target)
        return target

    reason = "missing"
    if target.is_file():
        if _read_installed_version(target) != EXPECTED_TEMPLATE_VERSION:
            reason = f"version {_read_installed_version(target)!r} != {EXPECTED_TEMPLATE_VERSION}"
        elif not template_structure_valid(target):
            reason = "structure (old Merci deck or missing Backlinks)"
    logger.info("Report template outdated (%s) — regenerating", reason)

    target.parent.mkdir(parents=True, exist_ok=True)
    _rebuild_template(target)

    if not template_is_current(target):
        raise RuntimeError(
            f"Template at {target} is still invalid after rebuild. "
            "Run: python scripts/build_template.py --force",
        )
    logger.info("Report template ready (%s slides): %s", EXPECTED_SLIDE_COUNT, target)
    return target
