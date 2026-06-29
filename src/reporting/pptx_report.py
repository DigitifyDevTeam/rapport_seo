"""Fill the PowerPoint template with the report data package.

The template (``templates/seo_report_template.pptx``) contains text and
picture placeholders written with ``{{name}}`` syntax. This module walks
every slide and replaces:

- text placeholders by their string value;
- picture placeholders by the matching chart image;
- table placeholders by a dataframe rendered as a styled PowerPoint table.
"""

from __future__ import annotations

import logging
import shutil
import tempfile
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from src.reporting.organic_performance_table import add_organic_performance_table
from src.transform.organic_performance import OrganicPerformanceSlide

logger = logging.getLogger(__name__)

PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG = RGBColor(0x0F, 0x17, 0x2A)
ROW_ALT = RGBColor(0xF1, 0xF5, 0xF9)
ROW_TEXT = RGBColor(0x1A, 0x1F, 0x36)

# Fractions of total table width per column (must sum to 1.0).
_TABLE_COLUMN_WIDTHS: dict[str, tuple[float, ...]] = {
    "table_top_pages": (0.68, 0.16, 0.16),
}

# Inset applied on each side of chart/screenshot placeholders (EMU fractions).
_PICTURE_MARGIN_RATIO = 0.04


class ReportBuilder:
    """Render the report deck for a single client and period."""

    def __init__(self, template_path: Path) -> None:
        self.template_path = template_path

    def _open_presentation(self):
        """Load the template; fall back to a temp copy if OneDrive locks the file."""
        try:
            return Presentation(self.template_path)
        except PermissionError:
            tmp = Path(tempfile.gettempdir()) / f"seo_report_template_{self.template_path.name}"
            try:
                shutil.copyfile(self.template_path, tmp)
            except PermissionError as exc:
                raise RuntimeError(
                    "Cannot read the PowerPoint template (permission denied). "
                    "Close apps locking the file, set OneDrive file to "
                    "'Always keep on this device', or set SEO_REPORT_TEMPLATE_PATH "
                    "in .env to a copy outside OneDrive."
                ) from exc
            logger.info("Opened PPTX template from temp copy %s (original locked)",
                        tmp)
            return Presentation(tmp)

    def build(self, data: dict[str, Any], output_path: Path) -> Path:
        prs = self._open_presentation()
        if data.get("clarity_hide_popular_products"):
            for slide in prs.slides:
                if self._slide_has_clarity_popular_products(slide):
                    self._remove_clarity_popular_products_column(slide)
        for slide in prs.slides:
            self._fill_slide(slide, data)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            prs.save(output_path)
            return output_path
        except PermissionError:
            from datetime import datetime
            stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
            fallback = output_path.with_name(
                f"{output_path.stem}.{stamp}{output_path.suffix}")
            logger.warning(
                "Could not overwrite %s (file open?). Saving to %s instead.",
                output_path, fallback,
            )
            prs.save(fallback)
            return fallback

    @staticmethod
    def _slide_has_clarity_popular_products(slide) -> bool:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text or ""
            if "chart_clarity_popular_products" in text:
                return True
        return False

    @staticmethod
    def _remove_clarity_popular_products_column(slide) -> None:
        """Drop Produits populaires and widen the remaining chart slot(s)."""
        clusters: list[list[Any]] = []
        tolerance = 120_000

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text or ""
            if (
                "chart_clarity_popular_products" in text
                or text.strip() == "Produits populaires"
            ):
                slide.shapes._spTree.remove(shape._element)  # noqa: SLF001
                continue
            if "chart_clarity_" not in text and text.strip() not in (
                "Appareils", "Référents", "Pages supérieures",
            ):
                continue
            placed = False
            for cluster in clusters:
                if abs(cluster[0].left - shape.left) <= tolerance:
                    cluster.append(shape)
                    placed = True
                    break
            if not placed:
                clusters.append([shape])

        if not clusters:
            return
        clusters.sort(key=lambda group: group[0].left)
        if len(clusters) == 1:
            group = clusters[0]
            anchor = group[0]
            panel_left = anchor.left
            panel_right = anchor.left + anchor.width
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text or ""
                if "chart_clarity_" not in text:
                    continue
                right = shape.left + shape.width
                if right > panel_right:
                    panel_right = right
                if shape.left < panel_left:
                    panel_left = shape.left
            span = panel_right - panel_left
            delta = panel_left - anchor.left
            for shape in group:
                shape.left = int(shape.left - delta)
                if shape.width > 0:
                    shape.width = span
            return
        gap = clusters[1][0].left - (clusters[0][0].left + clusters[0][0].width)
        if gap < 0:
            gap = 0
        left_edge = clusters[0][0].left
        right_edge = clusters[-1][0].left + clusters[-1][0].width
        span = right_edge - left_edge
        n = len(clusters)
        total_gap = gap * (n - 1)
        slot_w = int((span - total_gap) / n)
        for idx, group in enumerate(clusters):
            new_left = int(left_edge + idx * (slot_w + gap))
            old_left = group[0].left
            delta = new_left - old_left
            for shape in group:
                shape.left = int(shape.left + delta)
                if shape.width > 0:
                    shape.width = slot_w

    def _fill_slide(self, slide, data: dict[str, Any]) -> None:
        for shape in tuple(slide.shapes):
            if not shape.has_text_frame:
                continue
            placeholders = self._extract_placeholders(shape.text_frame.text)
            if not placeholders:
                continue
            self._handle_shape(slide, shape, placeholders, data)

    @staticmethod
    def _extract_placeholders(text: str) -> list[str]:
        names: list[str] = []
        cursor = 0
        while True:
            start = text.find("{{", cursor)
            if start < 0:
                break
            end = text.find("}}", start + 2)
            if end < 0:
                break
            names.append(text[start + 2:end].strip())
            cursor = end + 2
        return names

    def _handle_shape(self, slide, shape, placeholders: list[str],
                       data: dict[str, Any]) -> None:
        if len(placeholders) == 1:
            name = placeholders[0]
            if name.startswith("chart_"):
                self._replace_with_image(slide, shape, data.get(name),
                                         placeholder_name=name)
                return
            if name == "table_organic_performance":
                self._replace_organic_performance(slide, shape, data.get(name))
                return
            if name.startswith("table_"):
                self._replace_with_table(slide, shape, data.get(name),
                                          table_name=name)
                return
        if (len(placeholders) == 1
                and placeholders[0].startswith("final_summary_")):
            self._replace_final_summary(shape, placeholders[0], data)
            return
        self._replace_text(shape, placeholders, data)

    def _replace_final_summary(self, shape, name: str,
                                data: dict[str, Any]) -> None:
        body = _stringify(data.get(name))
        if name == "final_summary_brief":
            mode = "brief"
        elif name == "final_summary_takeaways":
            mode = "takeaway"
        else:
            mode = "card"
        self._set_summary_body(shape.text_frame, body, mode=mode)

    def _replace_text(self, shape, placeholders: list[str],
                       data: dict[str, Any]) -> None:
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                for name in placeholders:
                    token = "{{" + name + "}}"
                    if token in run.text:
                        run.text = run.text.replace(token,
                                                      _stringify(data.get(name)))

        full_text = "\n".join(p.text for p in tf.paragraphs)
        if any("{{" + p + "}}" in full_text for p in placeholders):
            new_text = full_text
            for name in placeholders:
                new_text = new_text.replace("{{" + name + "}}",
                                              _stringify(data.get(name)))
            self._set_text_frame(tf, new_text)

    @staticmethod
    def _set_summary_body(tf, text: str, *,
                          mode: str = "card") -> None:
        """Format Synthèse finale panel text.

        *mode*: "brief" (full-width intro), "card" (topic column),
                "takeaway" (dark strip at bottom).
        """
        sizes = {"brief": Pt(12), "card": Pt(11), "takeaway": Pt(10)}
        spacings = {"brief": Pt(5), "card": Pt(3), "takeaway": Pt(3)}
        max_lines = {"brief": 3, "card": 4, "takeaway": 6}
        base_size = sizes.get(mode, Pt(9))
        space = spacings.get(mode, Pt(2))
        cap = max_lines.get(mode, 5)

        for paragraph in list(tf.paragraphs)[1:]:
            paragraph._p.getparent().remove(paragraph._p)  # noqa: SLF001
        first = tf.paragraphs[0]
        first.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)

        lines = (text.split("\n") if text else [""])[:cap]
        is_dark = mode == "takeaway"
        for idx, line in enumerate(lines):
            paragraph = first if idx == 0 else tf.add_paragraph()
            paragraph.space_after = space
            paragraph.line_spacing = 1.08
            run = paragraph.add_run()
            run.text = line
            run.font.size = base_size
            run.font.name = "Segoe UI"
            if is_dark:
                run.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            elif line.strip().startswith("•"):
                run.font.color.rgb = ROW_TEXT
            else:
                run.font.color.rgb = PRIMARY

    @staticmethod
    def _set_text_frame(tf, text: str) -> None:
        first = tf.paragraphs[0]
        font = None
        if first.runs:
            font = first.runs[0].font
        for paragraph in list(tf.paragraphs)[1:]:
            paragraph._p.getparent().remove(paragraph._p)  # noqa: SLF001
        first.clear()
        tf.word_wrap = True
        lines = text.split("\n") or [""]
        for idx, line in enumerate(lines):
            paragraph = first if idx == 0 else tf.add_paragraph()
            paragraph.space_after = Pt(3)
            run = paragraph.add_run()
            run.text = line
            if font is not None and font.size is not None:
                run.font.size = font.size
            if font is not None and font.color and font.color.type is not None:
                try:
                    run.font.color.rgb = font.color.rgb
                except AttributeError:
                    pass
            stripped = line.strip()
            if stripped.isupper() and len(stripped) > 3:
                run.font.bold = True
                run.font.color.rgb = PRIMARY
            elif stripped.startswith("•"):
                run.font.color.rgb = ROW_TEXT
            elif stripped.startswith("À retenir"):
                run.font.bold = True

    def _replace_organic_performance(self, slide, shape, payload: Any) -> None:
        if not isinstance(payload, OrganicPerformanceSlide):
            self._set_text_frame(shape.text_frame, "Données GA4 indisponibles")
            return
        left, top, width, height = (shape.left, shape.top, shape.width,
                                      shape.height)
        slide.shapes._spTree.remove(shape._element)  # noqa: SLF001
        add_organic_performance_table(slide, left, top, width, height, payload)

    def _replace_with_image(self, slide, shape, image_path: Any, *,
                            placeholder_name: str = "") -> None:
        if not image_path:
            return
        path = Path(image_path)
        if not path.exists():
            logger.warning("chart image not found: %s", path)
            return
        left, top, width, height = (shape.left, shape.top, shape.width,
                                      shape.height)
        slide.shapes._spTree.remove(shape._element)  # noqa: SLF001
        fit_left, fit_top, fit_w, fit_h = _fitted_picture_bounds(
            left, top, width, height, path,
            margin_ratio=(
                0.012 if placeholder_name.startswith("chart_clarity_")
                else _PICTURE_MARGIN_RATIO
            ),
        )
        picture_path = path
        enhanced_path: Path | None = None
        if placeholder_name.startswith(("chart_gmb_", "chart_clarity_")):
            try:
                from src.reporting.screenshot_enhance import prepare_slide_image

                enhanced_path = prepare_slide_image(
                    path, fit_w, fit_h, placeholder_name=placeholder_name,
                )
                picture_path = enhanced_path
            except Exception as exc:
                logger.warning("screenshot enhance failed for %s: %s", path, exc)
        try:
            slide.shapes.add_picture(str(picture_path), fit_left, fit_top,
                                       width=fit_w, height=fit_h)
        finally:
            if enhanced_path is not None and enhanced_path != path:
                try:
                    enhanced_path.unlink(missing_ok=True)
                except OSError:
                    pass

    def _replace_with_table(self, slide, shape, df: Any,
                             table_name: str | None = None) -> None:
        if not isinstance(df, pd.DataFrame) or df.empty:
            self._set_text_frame(shape.text_frame, "No data available")
            return
        left, top, width, height = (shape.left, shape.top, shape.width,
                                      shape.height)
        slide.shapes._spTree.remove(shape._element)  # noqa: SLF001
        rows = len(df) + 1
        cols = len(df.columns)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width,
                                                height)
        table = table_shape.table
        if table_name:
            self._apply_column_widths(table, width, table_name)

        for col_idx, column in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(column)
            self._format_cell(cell, header=True)

        for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
            for col_idx, column in enumerate(df.columns):
                cell = table.cell(row_idx, col_idx)
                cell.text = _stringify(row[column])
                self._format_cell(cell, header=False,
                                    alt=row_idx % 2 == 0)

    @staticmethod
    def _apply_column_widths(table, total_width_emu: int,
                              table_name: str) -> None:
        fractions = _TABLE_COLUMN_WIDTHS.get(table_name)
        if not fractions or len(fractions) != len(table.columns):
            return
        scale = sum(fractions) or 1.0
        for idx, frac in enumerate(fractions):
            table.columns[idx].width = int(total_width_emu * frac / scale)

    @staticmethod
    def _format_cell(cell, *, header: bool, alt: bool = False) -> None:
        cell.fill.solid()
        if header:
            cell.fill.fore_color.rgb = HEADER_BG
        else:
            cell.fill.fore_color.rgb = (ROW_ALT if alt
                                          else RGBColor(0xFF, 0xFF, 0xFF))
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.bold = header
                run.font.color.rgb = (HEADER_TEXT if header else ROW_TEXT)


def _image_pixel_size(path: Path) -> tuple[int, int] | None:
    try:
        from PIL import Image

        with Image.open(path) as im:
            return im.size
    except Exception:
        return None


def _fitted_picture_bounds(left: int, top: int, width: int, height: int,
                            image_path: Path, *,
                            margin_ratio: float) -> tuple[int, int, int, int]:
    """Return EMU bounds that keep the image inside the placeholder with margin."""
    margin_x = int(width * margin_ratio)
    margin_y = int(height * margin_ratio)
    inner_w = max(width - 2 * margin_x, 1)
    inner_h = max(height - 2 * margin_y, 1)
    inner_left = left + margin_x
    inner_top = top + margin_y

    size = _image_pixel_size(image_path)
    if not size or size[0] <= 0 or size[1] <= 0:
        return inner_left, inner_top, inner_w, inner_h

    img_aspect = size[0] / size[1]
    box_aspect = inner_w / inner_h
    if img_aspect > box_aspect:
        pic_w = inner_w
        pic_h = max(int(inner_w / img_aspect), 1)
    else:
        pic_h = inner_h
        pic_w = max(int(inner_h * img_aspect), 1)

    pic_left = inner_left + (inner_w - pic_w) // 2
    pic_top = inner_top + (inner_h - pic_h) // 2
    return pic_left, pic_top, pic_w, pic_h


def _stringify(value: Any) -> str:
    if value is None:
        return "n/a"
    if isinstance(value, list):
        return "\n".join(f"\u2022 {_stringify(item)}" for item in value)
    if isinstance(value, float):
        if value.is_integer():
            return f"{int(value):,}"
        return f"{value:,.2f}"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)


def render(template_path: Path, output_path: Path,
            data: dict[str, Any]) -> Path:
    """Convenience wrapper for one-shot rendering."""
    return ReportBuilder(template_path).build(data, output_path)
