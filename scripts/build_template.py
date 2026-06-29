"""Generate the reusable PowerPoint template for the monthly SEO report.

The template is created programmatically so it can be versioned in git and
rebuilt on any machine without shipping a binary asset that drifts from the
code that fills it.

Run:
    python scripts/build_template.py          # first-time only (no file yet)
    python scripts/build_template.py --force  # overwrite an existing template

By default this script **does not** overwrite
``templates/seo_report_template.pptx`` if it already exists. Edit that file
in PowerPoint for layout/branding changes; monthly reports only fill
``{{placeholders}}`` and never regenerate the deck layout.

The output is written to ``templates/seo_report_template.pptx``.
"""

from __future__ import annotations

import argparse
import io
import os
import sys
from pathlib import Path

from dotenv import load_dotenv

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE_PATH = PROJECT_ROOT / "templates" / "seo_report_template.pptx"
# Bump when slide order/structure changes (keep in sync with ensure_template.py).
TEMPLATE_BUILD_VERSION = "2026-06-v5-clarity-hero-charts"


def resolve_template_path(output: str | None = None) -> Path:
    """Target .pptx: ``--output``, then ``SEO_REPORT_TEMPLATE_PATH``, then default."""
    if output:
        return Path(output).expanduser().resolve()
    load_dotenv(PROJECT_ROOT / ".env")
    override = (os.environ.get("SEO_REPORT_TEMPLATE_PATH") or "").strip()
    if override:
        return Path(override).expanduser().resolve()
    return DEFAULT_TEMPLATE_PATH


def template_version_path(pptx_path: Path) -> Path:
    return pptx_path.parent / ".template_build_version"

# Modern analytics deck: slate base + teal + violet (professional contrast)
PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x14, 0xB8, 0xA6)
ACCENT_BRIGHT = RGBColor(0x2D, 0xD4, 0xBF)
ACCENT_DIM = RGBColor(0x0F, 0x76, 0x6E)
ACCENT_SECOND = RGBColor(0x8B, 0x5C, 0xF6)
MUTED = RGBColor(0x64, 0x74, 0x8B)
TEXT = RGBColor(0x0F, 0x17, 0x2A)
LIGHT_BG = RGBColor(0xEC, 0xFE, 0xFF)
PAGE_BG = RGBColor(0xF8, 0xFA, 0xFC)
COVER_BG = RGBColor(0x04, 0x09, 0x12)
COVER_GLOW = RGBColor(0x0C, 0x16, 0x26)
COVER_MUTED = RGBColor(0x94, 0xA3, 0xB8)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
KPI_CARD_BORDER = RGBColor(0xAC, 0x7C, 0x78)
CARD_SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
KPI_ACCENT_PALETTE: list[RGBColor] = [ACCENT, ACCENT_SECOND, ACCENT_DIM, ACCENT_BRIGHT]
KPI_PILL_BG: dict[RGBColor, RGBColor] = {
    ACCENT: RGBColor(0xCC, 0xFB, 0xF1),
    ACCENT_SECOND: RGBColor(0xE9, 0xD5, 0xFF),
    ACCENT_DIM: RGBColor(0x99, 0xF6, 0xE4),
    ACCENT_BRIGHT: RGBColor(0xA7, 0xF3, 0xD0),
}

# Short client-facing definitions on the KPI overview slide.
KPI_PREAMBLES: dict[str, str] = {
    "Sessions": "Nombre de visites sur votre site.",
    "Utilisateurs": "Visiteurs uniques sur la période.",
    "Conversions": "Achats, contacts ou demandes importantes.",
    "Clics": "Clics depuis Google vers votre site.",
    "Impressions": "Fois où votre site apparaît sur Google.",
    "CTR": "Part des impressions qui devient un clic.",
    "Position moyenne": "Rang moyen sur Google (plus bas = mieux).",
}

# Slide numbers for the table of contents (cover = 1, ToC = 2, then content).
# Slide indices: 1 = cover, 2 = table of contents; keep in sync with ``main()`` order.
TOC_ITEMS: list[tuple[str, int]] = [
    ("Vue d'ensemble des KPI", 3),
    ("Performance organique (GA4)", 4),
    ("Trafic organique (GA4)", 5),
    ("Pages et écrans (GA4)", 6),
    ("Comportement (Clarity)", 7),
    ("Comportement (Clarity) — suite", 8),
    ("Performance Search (GSC)", 9),
    ("Top pages (GSC)", 10),
    ("Présence Google Business Profile", 11),
    ("Interactions clients (détail)", 12),
    ("Backlinks", 13),
    ("Synthèse finale", 15),
]

FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"

# Extra inset for screenshot slots (GMB, Clarity) inside their frames.
CHART_SLOT_INSET = Inches(0.1)

# Safe content zone (widescreen 13.333" × 7.5") — keeps elements inside the slide.
MARGIN_X = Inches(0.58)
MARGIN_BOTTOM = Inches(0.45)
CONTENT_TOP = Inches(1.20)
PANEL_INNER_PAD = Inches(0.32)
GRID_GAP = Inches(0.20)
ORGANIC_CONTENT_TOP = Inches(0.98)


def _set_text(shape, text: str, *, size: int = 18, bold: bool = False,
              color: RGBColor = TEXT, align=PP_ALIGN.LEFT,
              font_name: str = FONT_BODY) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_text_box(slide, left, top, width, height, text, **kwargs):
    box = slide.shapes.add_textbox(left, top, width, height)
    _set_text(box, text, **kwargs)
    return box


def _add_band(slide, left, top, width, height, color: RGBColor) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()


def _content_rect(prs: Presentation) -> tuple:
    """Main content box below the slide header (EMU)."""
    width = prs.slide_width - 2 * MARGIN_X
    height = prs.slide_height - CONTENT_TOP - MARGIN_BOTTOM
    return MARGIN_X, CONTENT_TOP, width, height


def _inner_rect(panel_left, panel_top, panel_w, panel_h) -> tuple:
    """Padding inside a white content panel (EMU)."""
    pad = PANEL_INNER_PAD
    return (panel_left + pad, panel_top + pad,
            panel_w - 2 * pad, panel_h - 2 * pad)


def _add_content_panel(slide, prs: Presentation,
                       *, fill: RGBColor = CARD_SURFACE) -> tuple:
    """Rounded panel filling the safe content area."""
    left, top, width, height = _content_rect(prs)
    _add_soft_panel(slide, left, top, width, height, fill=fill, line=CARD_BORDER)
    return left, top, width, height


def _fit_row(count: int, area_w, gap) -> int:
    """Cell width for a single horizontal row inside *area_w*."""
    if count < 1:
        return area_w
    return int((area_w - gap * (count - 1)) / count)


def _fit_grid(count: int, cols: int, area_w, area_h, gap_x, gap_y) -> tuple[int, int]:
    """Cell size for a grid that fits inside the given area."""
    rows = (count + cols - 1) // cols
    cell_w = int((area_w - gap_x * (cols - 1)) / cols)
    cell_h = int((area_h - gap_y * (rows - 1)) / rows)
    return cell_w, cell_h


def _organic_content_rect(prs: Presentation) -> tuple:
    """Content area on the organic performance slide (no title band)."""
    width = prs.slide_width - 2 * MARGIN_X
    height = prs.slide_height - ORGANIC_CONTENT_TOP - MARGIN_BOTTOM
    return MARGIN_X, ORGANIC_CONTENT_TOP, width, height


def _add_soft_panel(slide, left, top, width, height,
                    *, fill: RGBColor, line: RGBColor | None = None) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    if line is not None:
        panel.line.color.rgb = line
        panel.line.width = Pt(0.75)
    else:
        panel.line.fill.background()


def _kpi_value_font_size(height) -> int:
    """Scale headline KPI value to card height."""
    h_in = height / Inches(1)
    if h_in >= 2.2:
        return 32
    if h_in >= 1.45:
        return 26
    if h_in >= 1.05:
        return 22
    return 18


def _add_kpi_card(slide, left, top, width, height, label: str,
                   value_placeholder: str, delta_placeholder: str,
                   *, accent: RGBColor | None = None,
                   variant_index: int = 0,
                   compact: bool = False,
                   preamble: str | None = None) -> None:
    """KPI tile: terracotta outline, optional préambule, and delta badge."""
    accent_rgb = accent or KPI_ACCENT_PALETTE[variant_index % len(KPI_ACCENT_PALETTE)]
    pill_bg = KPI_PILL_BG.get(accent_rgb, LIGHT_BG)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    card.adjustments[0] = 0.14
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_SURFACE
    card.line.color.rgb = KPI_CARD_BORDER
    card.line.width = Pt(0.75)

    pad_x = Inches(0.14) if compact else Inches(0.2)
    pad_y = Inches(0.10) if compact else Inches(0.16)
    inner_left = left + pad_x
    inner_w = width - 2 * pad_x
    label_top = top + pad_y
    label_h = Inches(0.20) if compact else Inches(0.26)
    label_size = 7 if compact else 8
    _add_text_box(slide, inner_left, label_top, inner_w, label_h,
                  label.upper(), size=label_size, bold=True, color=MUTED)

    block_top = label_top + label_h
    if preamble:
        preamble_h = Inches(0.30) if compact else Inches(0.34)
        preamble_size = 7 if compact else 8
        _add_text_box(slide, inner_left, block_top + Inches(0.03),
                      inner_w, preamble_h, preamble,
                      size=preamble_size, color=MUTED)
        block_top += preamble_h + Inches(0.04)

    value_top = block_top + Inches(0.02 if compact else 0.04)
    value_size = _kpi_value_font_size(height)
    if compact:
        value_size = min(value_size, 20)
    value_h = Inches(0.38) + (value_size - 18) * Inches(0.015)
    _add_text_box(slide, inner_left, value_top, inner_w, value_h,
                  value_placeholder, size=value_size, bold=True,
                  color=PRIMARY, font_name=FONT_TITLE)

    if delta_placeholder:
        pill_h = Inches(0.3)
        pill_w = min(inner_w, Inches(1.45))
        pill_top = top + height - pad_y - pill_h
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      inner_left, pill_top, pill_w, pill_h)
        pill.adjustments[0] = 0.5
        pill.fill.solid()
        pill.fill.fore_color.rgb = pill_bg
        pill.line.fill.background()
        _add_text_box(slide, inner_left + Inches(0.1), pill_top + Inches(0.04),
                      pill_w - Inches(0.2), pill_h - Inches(0.06),
                      delta_placeholder, size=9, bold=True, color=accent_rgb)


def _slide_with_title(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height,
              PAGE_BG)
    rail_w = Inches(0.12)
    _add_band(slide, Inches(0), Inches(0), rail_w, prs.slide_height,
              ACCENT_SECOND)
    _add_band(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.048), ACCENT)
    _add_band(slide, Inches(0), Inches(0.048), prs.slide_width, Inches(0.024),
              ACCENT_BRIGHT)
    header_top = Inches(0.072)
    header_h = Inches(0.9)
    _add_band(slide, Inches(0), header_top, prs.slide_width, header_h, PRIMARY)
    title_left = rail_w + Inches(0.42)
    title_w = prs.slide_width - title_left - MARGIN_X
    _add_text_box(slide, title_left, header_top + Inches(0.1), title_w,
                  Inches(0.5), title, size=22, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF),
                  font_name=FONT_TITLE)
    if subtitle:
        _add_text_box(slide, title_left, header_top + Inches(0.56), title_w,
                      Inches(0.36), subtitle, size=12,
                      color=RGBColor(0xBA, 0xCC, 0xE8))
    _add_band(slide, Inches(0), header_top + header_h,
              prs.slide_width, Inches(0.038), ACCENT_DIM)
    return slide


def _add_toc_entry(slide, left: float, top: float, width: float, page: int,
                    title: str, badge_color: RGBColor) -> None:
    badge_w = Inches(0.42)
    badge_h = Inches(0.36)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, badge_w, badge_h)
    badge.adjustments[0] = 0.35
    badge.fill.solid()
    badge.fill.fore_color.rgb = badge_color
    badge.line.fill.background()
    page_str = str(page)
    _add_text_box(slide, left, top + Inches(0.02), badge_w, badge_h - Inches(0.02),
                  page_str, size=12, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    text_left = left + badge_w + Inches(0.16)
    text_w = width - badge_w - Inches(0.16)
    _add_text_box(slide, text_left, top + Inches(0.04), text_w, badge_h,
                  title, size=13, color=TEXT, bold=False)


def build_table_of_contents(prs: Presentation) -> None:
    slide = _slide_with_title(prs, "Table des matières",
                               "Plan du rapport · navigation rapide")
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)
    col_w = int((inner_w - GRID_GAP) / 2)
    row_h = Inches(0.44)
    start_top = inner_top + Inches(0.06)
    left_col_x = inner_left
    right_col_x = inner_left + col_w + GRID_GAP
    half = (len(TOC_ITEMS) + 1) // 2
    for idx, (toc_title, page_num) in enumerate(TOC_ITEMS):
        col = 0 if idx < half else 1
        row = idx if idx < half else idx - half
        left = left_col_x if col == 0 else right_col_x
        top = start_top + row * row_h
        badge_color = ACCENT if idx % 2 == 0 else ACCENT_SECOND
        _add_toc_entry(slide, left, top, col_w, page_num, toc_title,
                        badge_color)


def _picture_placeholder(slide, left, top, width, height, name: str) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = f"{{{{{name}}}}}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED
        run.font.italic = True


def _add_framed_picture_placeholder(slide, left, top, width, height,
                                   name: str, *,
                                   inset=None) -> None:
    """Chart/screenshot slot with a visible rounded border (GMB detail grid)."""
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    frame.adjustments[0] = 0.06
    frame.fill.solid()
    frame.fill.fore_color.rgb = CARD_SURFACE
    frame.line.color.rgb = CARD_BORDER
    frame.line.width = Pt(0.75)
    slot_inset = CHART_SLOT_INSET if inset is None else inset
    _picture_placeholder(
        slide,
        left + slot_inset,
        top + slot_inset,
        width - 2 * slot_inset,
        height - 2 * slot_inset,
        name,
    )


def _table_placeholder(slide, left, top, width, height, name: str,
                        caption: str) -> None:
    _add_text_box(slide, left, top, width, Inches(0.3),
                  caption, size=12, bold=True, color=PRIMARY)
    _picture_placeholder(slide, left, top + Inches(0.3), width,
                          height - Inches(0.3), name)


_COVER_PROFILE_ROWS: list[tuple[str, str]] = [
    ("Client", "{{cover_client}}"),
    ("Activité", "{{cover_activity}}"),
    ("Nom du site", "{{cover_site_name}}"),
    ("URL", "{{cover_url}}"),
    ("Pack SEO", "{{cover_seo_pack}}"),
    ("Première activité SEO", "{{cover_seo_since}}"),
]


def _add_cover_profile_panel(slide, left, top, width, _height) -> None:
    """Project facts inside the cover slide dark panel."""
    pad_x = Inches(0.4)
    pad_y = Inches(0.38)
    inner_left = left + pad_x
    inner_w = width - 2 * pad_x
    divider_rgb = RGBColor(0x22, 0x33, 0x55)
    label_rgb = RGBColor(0x94, 0xA3, 0xB8)
    value_rgb = RGBColor(0xF1, 0xF5, 0xF9)

    title_top = top + pad_y
    title_h = Inches(0.26)
    _add_text_box(slide, inner_left, title_top, inner_w, title_h,
                  "FICHE PROJET", size=9, bold=True, color=ACCENT_BRIGHT,
                  font_name=FONT_TITLE)
    _add_band(slide, inner_left, title_top + title_h + Inches(0.08), inner_w,
              Inches(0.045), ACCENT)

    row_top = title_top + title_h + Inches(0.28)
    label_h = Inches(0.2)
    value_h = Inches(0.34)
    row_gap = Inches(0.14)
    divider_h = Inches(0.018)

    for idx, (label, value_placeholder) in enumerate(_COVER_PROFILE_ROWS):
        _add_text_box(slide, inner_left, row_top, inner_w, label_h,
                      label.upper(), size=8, bold=True, color=label_rgb)
        _add_text_box(slide, inner_left, row_top + label_h, inner_w, value_h,
                      value_placeholder, size=12, bold=True, color=value_rgb,
                      font_name=FONT_TITLE)
        row_bottom = row_top + label_h + value_h + Inches(0.06)
        if idx < len(_COVER_PROFILE_ROWS) - 1:
            _add_band(slide, inner_left, row_bottom, inner_w, divider_h,
                      divider_rgb)
            row_top = row_bottom + divider_h + row_gap
        else:
            row_top = row_bottom


def build_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height,
              COVER_BG)
    _add_band(slide, Inches(0), Inches(0), Inches(0.14), prs.slide_height,
              ACCENT_SECOND)
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        prs.slide_width - Inches(3.8), Inches(-2.1), Inches(7.2), Inches(7.2))
    glow.fill.solid()
    glow.fill.fore_color.rgb = COVER_GLOW
    glow.line.fill.background()
    panel_left = prs.slide_width - Inches(5.15)
    panel_top = Inches(0.85)
    panel_w = Inches(4.85)
    panel_h = Inches(5.65)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    panel_left, panel_top, panel_w, panel_h)
    panel.adjustments[0] = 0.06
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x08, 0x11, 0x20)
    panel.line.color.rgb = RGBColor(0x22, 0x33, 0x55)
    panel.line.width = Pt(0.5)
    _add_cover_profile_panel(slide, panel_left, panel_top, panel_w, panel_h)
    _add_band(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.12), ACCENT)
    _add_band(slide, Inches(0), Inches(0.12), prs.slide_width, Inches(0.05),
              ACCENT_BRIGHT)
    # Bottom accent only on the left column — do not cross the Fiche projet panel.
    _add_band(slide, Inches(0), Inches(5.72), panel_left, Inches(0.07),
              ACCENT_DIM)
    _add_band(slide, prs.slide_width - Inches(0.28), Inches(0.12),
              Inches(0.28), Inches(7.38), ACCENT_SECOND)
    _add_text_box(slide, Inches(0.62), Inches(1.22), Inches(10.8), Inches(0.45),
                  "{{agency_name}}", size=11, color=COVER_MUTED, bold=True)
    _add_text_box(slide, Inches(0.62), Inches(1.68), Inches(10.8), Inches(1.05),
                  "Rapport SEO", size=44, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), font_name=FONT_TITLE)
    _add_text_box(slide, Inches(0.62), Inches(2.58), Inches(10.8), Inches(0.55),
                  "mensuel", size=26, bold=False,
                  color=ACCENT_BRIGHT, font_name=FONT_TITLE)
    _add_band(slide, Inches(0.62), Inches(3.32), Inches(1.38), Inches(0.052),
              ACCENT)
    _add_band(slide, Inches(2.0), Inches(3.32), Inches(1.38), Inches(0.052),
              ACCENT_SECOND)
    _add_text_box(slide, Inches(0.62), Inches(3.48), Inches(10.8), Inches(0.68),
                  "{{client_name}}", size=26, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), font_name=FONT_TITLE)
    _add_text_box(slide, Inches(0.62), Inches(4.18), Inches(10.8), Inches(0.45),
                  "{{period_label}}", size=15, color=COVER_MUTED)
    _add_text_box(slide, Inches(0.62), Inches(6.68), Inches(11.5), Inches(0.42),
                  "Document confidentiel · Généré le {{report_date}}",
                  size=10, color=RGBColor(0x6A, 0x78, 0x8C))


_SUMMARY_SECTION_COLORS: list[tuple[RGBColor, RGBColor, RGBColor]] = [
    # (icon_bg, accent_line, subtle_bg)
    (ACCENT, RGBColor(0xCC, 0xFB, 0xF1), RGBColor(0xF0, 0xFD, 0xFA)),
    (ACCENT_SECOND, RGBColor(0xE9, 0xD5, 0xFF), RGBColor(0xF5, 0xF3, 0xFF)),
    (ACCENT_DIM, RGBColor(0x99, 0xF6, 0xE4), RGBColor(0xF0, 0xFD, 0xFA)),
    (ACCENT_BRIGHT, RGBColor(0xA7, 0xF3, 0xD0), RGBColor(0xEC, 0xFD, 0xF5)),
]


def _add_summary_card(slide, left, top, width, height, title: str,
                      placeholder: str, *, color_idx: int = 0) -> None:
    """Topic card for Synthèse finale — icon badge + accent top bar + body."""
    icon_bg, pill_tint, card_bg = _SUMMARY_SECTION_COLORS[
        color_idx % len(_SUMMARY_SECTION_COLORS)]

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = card_bg
    card.line.color.rgb = pill_tint
    card.line.width = Pt(0.75)

    bar_h = Inches(0.04)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = icon_bg
    bar.line.fill.background()

    badge_size = Inches(0.28)
    badge_left = left + Inches(0.14)
    badge_top = top + bar_h + Inches(0.10)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   badge_left, badge_top,
                                   badge_size, badge_size)
    badge.adjustments[0] = 0.25
    badge.fill.solid()
    badge.fill.fore_color.rgb = icon_bg
    badge.line.fill.background()
    icon_char = title[0]
    _add_text_box(slide, badge_left, badge_top, badge_size, badge_size,
                  icon_char, size=12, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                  font_name=FONT_TITLE)

    title_left = badge_left + badge_size + Inches(0.10)
    title_w = width - (title_left - left) - Inches(0.14)
    title_h = Inches(0.26)
    _add_text_box(slide, title_left, badge_top, title_w, title_h,
                  title.upper(), size=10, bold=True, color=PRIMARY,
                  font_name=FONT_TITLE)

    body_top = badge_top + max(badge_size, title_h) + Inches(0.06)
    body_left = left + Inches(0.14)
    body_w = width - Inches(0.28)
    body_h = top + height - body_top - Inches(0.10)
    _add_text_box(slide, body_left, body_top, body_w, body_h,
                  f"{{{{{placeholder}}}}}", size=11, color=TEXT)


def _add_summary_highlight(slide, left, top, width, height, title: str,
                           placeholder: str, *, dark: bool = False) -> None:
    """Full-width highlight strip for brief / takeaways."""
    if dark:
        fill = PRIMARY
        title_rgb = ACCENT_BRIGHT
        body_rgb = RGBColor(0xE2, 0xE8, 0xF0)
        border = RGBColor(0x1E, 0x29, 0x3B)
    else:
        fill = RGBColor(0xF0, 0xFD, 0xFA)
        title_rgb = ACCENT_DIM
        body_rgb = PRIMARY
        border = ACCENT

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(0.75)

    accent_w = Inches(0.05)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 left, top, accent_w, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT if not dark else ACCENT_BRIGHT
    bar.line.fill.background()

    pad_x = Inches(0.18)
    pad_y = Inches(0.10)
    text_left = left + accent_w + pad_x
    text_w = width - accent_w - 2 * pad_x
    title_h = Inches(0.22)
    _add_text_box(slide, text_left, top + pad_y, text_w, title_h,
                  title.upper(), size=9, bold=True, color=title_rgb,
                  font_name=FONT_TITLE)
    body_top = top + pad_y + title_h + Inches(0.02)
    body_h = height - pad_y - title_h - Inches(0.02) - pad_y
    _add_text_box(slide, text_left, body_top, text_w, body_h,
                  f"{{{{{placeholder}}}}}", size=12, color=body_rgb)


def build_backlinks_slide(prs: Presentation, *, part: int = 1) -> None:
    """Empty Backlinks slide — content added manually in PowerPoint."""
    subtitle = "À compléter" if part == 1 else "À compléter (suite)"
    slide = _slide_with_title(prs, "Backlinks", subtitle)
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)
    _add_soft_panel(slide, inner_left, inner_top, inner_w, inner_h,
                    fill=RGBColor(0xFF, 0xFF, 0xFF), line=CARD_BORDER)


def build_final_summary_slide(prs: Presentation) -> None:
    """Synthèse finale — brief banner + 4 topic cards."""
    slide = _slide_with_title(prs, "Synthèse finale",
                                "Performance et recommandations")
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)

    gap = Inches(0.14)
    brief_h = Inches(0.72)
    grid_h = inner_h - brief_h - 2 * gap
    cols = 4
    col_w = int((inner_w - gap * (cols - 1)) / cols)

    y = inner_top
    _add_summary_highlight(slide, inner_left, y, inner_w, brief_h,
                           "En bref", "final_summary_brief")
    y += brief_h + gap

    cards = [
        ("Site web", "final_summary_website"),
        ("Visibilité Google", "final_summary_search"),
        ("Exp. visiteurs", "final_summary_clarity"),
        ("Fiche Google", "final_summary_gmb"),
    ]
    for idx, (title, ph) in enumerate(cards):
        x = inner_left + idx * (col_w + gap)
        _add_summary_card(slide, x, y, col_w, grid_h,
                          title, ph, color_idx=idx)


def build_kpi_overview(prs: Presentation) -> None:
    slide = _slide_with_title(prs, "Vue d'ensemble des KPI",
                                "Performance mois sur mois")
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)
    cards = [
        ("Sessions", "{{sessions}}", "{{sessions_delta}}"),
        ("Utilisateurs", "{{users}}", "{{users_delta}}"),
        ("Conversions", "{{conversions}}", "{{conversions_delta}}"),
        ("Clics", "{{clicks}}", "{{clicks_delta}}"),
        ("Impressions", "{{impressions}}", "{{impressions_delta}}"),
        ("CTR", "{{ctr}}", "{{ctr_delta}}"),
        ("Position moyenne", "{{avg_position}}", "{{avg_position_delta}}"),
    ]
    cols = 4
    gap = GRID_GAP
    card_w, card_h = _fit_grid(len(cards), cols, inner_w, inner_h, gap, gap)
    for idx, (label, value, delta) in enumerate(cards):
        row, col = divmod(idx, cols)
        left = inner_left + col * (card_w + gap)
        top = inner_top + row * (card_h + gap)
        _add_kpi_card(slide, left, top, card_w, card_h, label, value, delta,
                       variant_index=idx,
                       preamble=KPI_PREAMBLES.get(label))


def _add_chart_synthesis_panel(slide, left, top, width, height,
                                commentary_name: str) -> None:
    """Right-hand synthesis block on GA4 / GSC chart slides."""
    pad = Inches(0.14)
    inner_left = left + pad
    inner_w = width - 2 * pad
    _add_text_box(slide, inner_left, top + pad, inner_w, Inches(0.34),
                  "Synthèse", size=13, bold=True, color=PRIMARY,
                  font_name=FONT_TITLE)
    _add_band(slide, inner_left, top + pad + Inches(0.36), Inches(1.05),
              Inches(0.045), ACCENT)
    _add_text_box(slide, inner_left, top + pad + Inches(0.48),
                  inner_w, height - pad - Inches(0.48),
                  f"{{{{{commentary_name}}}}}", size=10, color=TEXT)


# Organic performance slide uses the shared KPI card design.


def build_organic_performance(prs: Presentation) -> None:
    """GA4 organic channel summary (KPI row + period comparison table)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height,
              PAGE_BG)
    _add_band(slide, Inches(0), Inches(0), Inches(0.12), prs.slide_height,
              ACCENT_SECOND)
    _add_band(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.048), ACCENT)

    title_top = Inches(0.42)
    content_left, content_top, content_w, content_h = _organic_content_rect(prs)
    _add_text_box(slide, content_left, title_top, content_w, Inches(0.55),
                  "{{organic_performance_title}}", size=20, bold=True,
                  color=PRIMARY, font_name=FONT_TITLE)

    inner_left = content_left + PANEL_INNER_PAD
    inner_top = content_top + PANEL_INNER_PAD
    inner_w = content_w - 2 * PANEL_INNER_PAD
    inner_h = content_h - 2 * PANEL_INNER_PAD
    footnote_h = Inches(0.38)
    gap = GRID_GAP
    kpi_defs = [
        ("Utilisateurs*", "{{organic_perf_users}}"),
        ("Nouveaux utilisateurs*", "{{organic_perf_new_users}}"),
        ("Sessions*", "{{organic_perf_sessions}}"),
        ("Taux d'engagement*", "{{organic_perf_engagement}}"),
    ]
    kpi_h = int(Inches(1.08))
    kpi_w = _fit_row(len(kpi_defs), inner_w, gap)
    for idx, (label, placeholder) in enumerate(kpi_defs):
        left = inner_left + idx * (kpi_w + gap)
        _add_kpi_card(slide, left, inner_top, kpi_w, kpi_h, label,
                       placeholder, "", variant_index=idx)

    table_top = inner_top + kpi_h + gap
    table_h = inner_h - kpi_h - gap - footnote_h
    _picture_placeholder(slide, inner_left, table_top, inner_w, table_h,
                          "table_organic_performance")

    _add_text_box(slide, inner_left,
                  content_top + content_h - PANEL_INNER_PAD - footnote_h,
                  inner_w, footnote_h,
                  "* : Visites et utilisateurs venant depuis les moteurs de "
                  "recherche seulement",
                  size=9, color=MUTED)


def build_chart_slide(prs: Presentation, title: str, subtitle: str,
                       chart_name: str, commentary_name: str) -> None:
    slide = _slide_with_title(prs, title, subtitle)
    panel_left, panel_top, panel_w, panel_h = _content_rect(prs)
    split_gap = Inches(0.22)
    synth_w = int(panel_w * 0.28)
    chart_w = panel_w - synth_w - split_gap
    _add_soft_panel(slide, panel_left, panel_top, chart_w, panel_h,
                    fill=RGBColor(0xFF, 0xFF, 0xFF), line=CARD_BORDER)
    synth_left = panel_left + chart_w + split_gap
    _add_soft_panel(slide, synth_left, panel_top, synth_w, panel_h,
                    fill=LIGHT_BG, line=CARD_BORDER)
    chart_pad = Inches(0.22)
    _picture_placeholder(
        slide,
        panel_left + chart_pad,
        panel_top + chart_pad,
        chart_w - 2 * chart_pad,
        panel_h - 2 * chart_pad,
        chart_name,
    )
    _add_chart_synthesis_panel(slide, synth_left, panel_top, synth_w, panel_h,
                                commentary_name)


def build_table_slide(prs: Presentation, title: str, subtitle: str,
                       table_name: str) -> None:
    slide = _slide_with_title(prs, title, subtitle)
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)
    _table_placeholder(slide, inner_left, inner_top, inner_w, inner_h,
                       table_name, subtitle)


def build_gmb_overview(prs: Presentation) -> None:
    """Slide 11: Knowledge Panel capture + five KPI placeholders (3 + 2 grid)."""
    slide = _slide_with_title(
        prs,
        "Présence Google Business Profile",
        "Fiche d'établissement et interactions clients",
    )
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)

    split_gap = Inches(0.22)
    card_w = int(inner_w * 0.47)
    card_left = inner_left
    card_top = inner_top
    card_h = inner_h
    _add_text_box(slide, card_left, card_top, card_w, Inches(0.3),
                   "Fiche publique", size=12, bold=True, color=PRIMARY,
                   align=PP_ALIGN.CENTER)
    card_img_top = card_top + Inches(0.3)
    card_img_h = card_h - Inches(0.3)
    _picture_placeholder(
        slide,
        card_left + CHART_SLOT_INSET,
        card_img_top + CHART_SLOT_INSET,
        card_w - 2 * CHART_SLOT_INSET,
        card_img_h - 2 * CHART_SLOT_INSET,
        "chart_gmb_business_card",
    )

    kpi_cards = [
        ("Interactions totales", "{{gmb_interactions}}"),
        ("Appels", "{{gmb_calls}}"),
        ("Réservations", "{{gmb_bookings}}"),
        ("Itinéraires", "{{gmb_directions}}"),
        ("Clics vers le site Web", "{{gmb_website_clicks}}"),
    ]
    kpi_left = card_left + card_w + split_gap
    kpi_area_w = inner_left + inner_w - kpi_left
    cols = 3
    gap_x = GRID_GAP
    gap_y = GRID_GAP
    kpi_w = int((kpi_area_w - gap_x * (cols - 1)) / cols)
    rows = 2
    # KPI height: between the original full split and the previous 0.92" cap.
    kpi_h_natural = int((card_h - gap_y * (rows - 1)) / rows)
    kpi_h = min(kpi_h_natural, int(Inches(1.22)))
    kpi_h = max(kpi_h, int(Inches(1.05)))
    kpi_grid_h = rows * kpi_h + (rows - 1) * gap_y
    kpi_block_top = card_top + max(0, (card_h - kpi_grid_h) // 2)
    for idx, (label, value) in enumerate(kpi_cards):
        row, col = divmod(idx, cols)
        if row == 1:
            row2_count = len(kpi_cards) - cols
            row2_total_w = kpi_w * row2_count + gap_x * (row2_count - 1)
            row2_start = kpi_left + (kpi_area_w - row2_total_w) / 2
            left = row2_start + (idx - cols) * (kpi_w + gap_x)
        else:
            left = kpi_left + col * (kpi_w + gap_x)
        top = kpi_block_top + row * (kpi_h + gap_y)
        _add_kpi_card(slide, left, top, kpi_w, kpi_h, label, value, "",
                       variant_index=idx)


def build_gmb_details(prs: Presentation) -> None:
    """Slide 12: five Performance tab screenshots (3 + 2 grid)."""
    slide = _slide_with_title(
        prs,
        "Interactions clients (détail)",
        "Performance par type d'interaction",
    )
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)

    charts = [
        ("Vue d'ensemble", "chart_gmb_overview"),
        ("Appels", "chart_gmb_calls"),
        ("Réservations", "chart_gmb_bookings"),
        ("Itinéraires", "chart_gmb_directions"),
        ("Clics vers le site Web", "chart_gmb_website_clicks"),
    ]
    cols = 3
    gap_x = GRID_GAP
    gap_y = Inches(0.18)
    caption_h = Inches(0.3)
    chart_w, chart_h = _fit_grid(len(charts), cols, inner_w, inner_h, gap_x, gap_y)
    chart_h = int(chart_h - caption_h)
    total_w = chart_w * cols + gap_x * (cols - 1)
    start_left = inner_left + int((inner_w - total_w) / 2)
    start_top = inner_top
    for idx, (caption, name) in enumerate(charts):
        row, col = divmod(idx, cols)
        if row == 1:
            row2_count = len(charts) - cols
            row2_w = chart_w * row2_count + gap_x * (row2_count - 1)
            r2_left = inner_left + int((inner_w - row2_w) / 2)
            left = r2_left + (idx - cols) * (chart_w + gap_x)
        else:
            left = start_left + col * (chart_w + gap_x)
        top = start_top + row * (chart_h + gap_y + caption_h)
        _add_text_box(slide, left, top, chart_w, caption_h,
                       caption, size=11, bold=True, color=PRIMARY,
                       align=PP_ALIGN.CENTER)
        img_top = top + caption_h
        img_h = chart_h
        _add_framed_picture_placeholder(
            slide, left, img_top, chart_w, img_h, name)


def _add_clarity_chart_row(
    slide,
    charts: list[tuple[str, str]],
    inner_left: int,
    charts_top: int,
    inner_w: int,
    charts_area_h: int,
    *,
    hero: bool = False,
) -> int:
    """Lay out *charts* in one row; return the Y coordinate below the row."""
    gap = Inches(0.16) if hero else GRID_GAP
    caption_h = Inches(0.24) if hero else Inches(0.3)
    chart_w = _fit_row(len(charts), inner_w, gap)
    chart_h = int(charts_area_h - caption_h)
    charts_total_w = chart_w * len(charts) + gap * (len(charts) - 1)
    charts_left = inner_left + int((inner_w - charts_total_w) / 2)
    for idx, (caption, name) in enumerate(charts):
        left = charts_left + idx * (chart_w + gap)
        _add_text_box(
            slide, left, charts_top, chart_w, caption_h,
            caption, size=12 if hero else 11, bold=True, color=PRIMARY,
            align=PP_ALIGN.LEFT if hero else PP_ALIGN.CENTER,
        )
        if hero:
            _add_band(slide, left, charts_top + caption_h - Inches(0.02),
                      Inches(0.72), Inches(0.035), ACCENT)
        img_top = charts_top + caption_h
        if hero:
            _add_framed_picture_placeholder(
                slide, left, img_top, chart_w, chart_h, name,
                inset=Inches(0.03),
            )
        else:
            _picture_placeholder(
                slide,
                left + CHART_SLOT_INSET,
                img_top + CHART_SLOT_INSET,
                chart_w - 2 * CHART_SLOT_INSET,
                chart_h - 2 * CHART_SLOT_INSET,
                name,
            )
    return charts_top + caption_h + chart_h


def build_clarity(prs: Presentation) -> None:
    gap = GRID_GAP
    commentary_h = Inches(0.38)

    slide1 = _slide_with_title(
        prs, "Comportement (Clarity)", "Signaux d'expérience utilisateur",
    )
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide1, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)

    kpi_cards = [
        ("Sessions", "{{clarity_sessions}}"),
        ("Pages par session", "{{clarity_pages_per_session}}"),
        ("Profondeur de défilement", "{{clarity_scroll_depth}}"),
        ("Temps d'activité passé", "{{clarity_active_time}}"),
    ]
    kpi_h = int(Inches(0.84))
    kpi_w = _fit_row(len(kpi_cards), inner_w, gap)
    cards_top = inner_top
    for idx, (label, value) in enumerate(kpi_cards):
        left = inner_left + idx * (kpi_w + gap)
        _add_kpi_card(
            slide1, left, cards_top, kpi_w, kpi_h, label, value, "",
            variant_index=idx, compact=True,
        )

    charts_gap = Inches(0.14)
    charts_top = cards_top + kpi_h + charts_gap
    charts_area_h = inner_h - kpi_h - charts_gap
    _add_clarity_chart_row(
        slide1,
        [
            ("Appareils", "chart_clarity_devices"),
            ("Référents", "chart_clarity_referrers"),
        ],
        inner_left,
        charts_top,
        inner_w,
        charts_area_h,
        hero=True,
    )

    slide2 = _slide_with_title(
        prs, "Comportement (Clarity)", "Pages et produits populaires",
    )
    panel_left, panel_top, panel_w, panel_h = _add_content_panel(slide2, prs)
    inner_left, inner_top, inner_w, inner_h = _inner_rect(
        panel_left, panel_top, panel_w, panel_h)
    charts_area_h = inner_h - commentary_h - gap
    charts_bottom = _add_clarity_chart_row(
        slide2,
        [
            ("Pages supérieures", "chart_clarity_popular_pages"),
            ("Produits populaires", "chart_clarity_popular_products"),
        ],
        inner_left,
        inner_top,
        inner_w,
        charts_area_h,
        hero=True,
    )
    _add_text_box(
        slide2, inner_left, charts_bottom + gap, inner_w, commentary_h,
        "{{clarity_commentary}}", size=11, color=TEXT, align=PP_ALIGN.CENTER,
    )


def _render_thank_you_gauge_png() -> bytes:
    """Semi-circular gauge (5 segments + needle) for the closing slide."""
    colors = ["#B2E5B5", "#1B5E20", "#FFCC80", "#FB8C00", "#D32F2F"]
    fig = plt.figure(figsize=(2.65, 1.38), dpi=160)
    ax = fig.add_subplot(projection="polar")
    ax.set_theta_zero_location("W")
    ax.set_theta_direction(-1)
    ax.set_axis_off()
    ax.set_ylim(0, 1.05)
    try:
        ax.set_thetamin(0)
        ax.set_thetamax(180)
    except (AttributeError, ValueError):
        pass
    width = np.pi / 5
    bottom = 0.18
    for i in range(5):
        theta_c = np.pi - (i + 0.5) * width
        ax.bar(
            theta_c,
            0.72,
            width=width * 0.92,
            bottom=bottom,
            color=colors[i],
            align="center",
            edgecolor="#222222",
            linewidth=0.35,
        )
    needle = np.pi - 1.5 * width
    ax.plot(
        [needle, needle],
        [0.05, bottom + 0.68],
        color="#111111",
        linewidth=2.8,
        zorder=10,
    )
    ax.scatter([needle], [0.06], s=60, color="#111111", zorder=11, edgecolors="none")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True, bbox_inches="tight",
                pad_inches=0.02)
    plt.close(fig)
    return buf.getvalue()


def build_thank_you_slide(prs: Presentation) -> None:
    """Closing slide: thank-you message, gauge, and agency contact (DIGITIFY)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    thank_bg = RGBColor(0xFA, 0xFA, 0xFA)
    _add_band(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height,
              thank_bg)

    box_w = Inches(8.85)
    box_h = Inches(2.88)
    box_left = (prs.slide_width - box_w) / 2
    box_top = Inches(2.18)

    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, box_top,
                                   box_w, box_h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    panel.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    panel.line.width = Pt(1.5)

    gauge_w = Inches(2.52)
    gauge_h = Inches(1.32)
    gx = prs.slide_width / 2 - gauge_w / 2
    gy = box_top - gauge_h + Inches(0.28)
    stream = io.BytesIO(_render_thank_you_gauge_png())
    slide.shapes.add_picture(stream, gx, gy, width=gauge_w, height=gauge_h)

    green = RGBColor(0x1B, 0x5E, 0x20)
    inner_pad = Inches(0.35)
    text_top = box_top + Inches(1.05)
    _add_text_box(slide, box_left + inner_pad, text_top,
                  box_w - 2 * inner_pad, Inches(0.68),
                  "MERCI POUR VOTRE", size=26, bold=True, color=green,
                  align=PP_ALIGN.CENTER, font_name=FONT_TITLE)
    _add_text_box(slide, box_left + inner_pad, text_top + Inches(0.62),
                  box_w - 2 * inner_pad, Inches(0.68),
                  "ATTENTION", size=26, bold=True, color=green,
                  align=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    contact_top = box_top + box_h + Inches(0.4)
    black = RGBColor(0x00, 0x00, 0x00)
    _add_text_box(slide, Inches(0.55), contact_top,
                  prs.slide_width - Inches(1.1), Inches(0.42),
                  "N'hésitez pas à nous contacter pour plus d'information",
                  size=12, color=black, align=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(0.55), contact_top + Inches(0.4),
                  prs.slide_width - Inches(1.1), Inches(0.36),
                  "contact@digitify.fr", size=12, bold=True, color=black,
                  align=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(0.55), contact_top + Inches(0.76),
                  prs.slide_width - Inches(1.1), Inches(0.36),
                  "07 45 80 14 12", size=12, bold=True, color=black,
                  align=PP_ALIGN.CENTER)


def _installed_template_version(pptx_path: Path) -> str | None:
    version_path = template_version_path(pptx_path)
    if not version_path.is_file():
        return None
    try:
        return version_path.read_text(encoding="utf-8").strip() or None
    except OSError:
        return None


def template_needs_rebuild(pptx_path: Path | None = None) -> bool:
    """True when the .pptx is missing or was built from older deck code."""
    target = pptx_path or resolve_template_path()
    if not target.is_file():
        return True
    return _installed_template_version(target) != TEMPLATE_BUILD_VERSION


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate the SEO report PowerPoint template from code.",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Overwrite templates/seo_report_template.pptx if it already exists.",
    )
    parser.add_argument(
        "--force-if-stale",
        action="store_true",
        help="Overwrite only when .template_build_version is outdated (used by cron).",
    )
    parser.add_argument(
        "--output",
        default="",
        help="Output .pptx path (default: SEO_REPORT_TEMPLATE_PATH or templates/).",
    )
    return parser.parse_args()


def _write_template_version(pptx_path: Path) -> None:
    template_version_path(pptx_path).write_text(
        f"{TEMPLATE_BUILD_VERSION}\n",
        encoding="utf-8",
    )


def _build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover(prs)
    build_table_of_contents(prs)
    build_kpi_overview(prs)
    build_organic_performance(prs)
    build_chart_slide(prs, "Trafic organique (GA4)",
                       "Sessions et utilisateurs sur la période",
                       "chart_ga4_traffic", "ga4_commentary")
    build_chart_slide(prs, "Pages et écrans (GA4)",
                       "Engagement — vues par jour sur la période",
                       "chart_ga4_pages_screens", "ga4_pages_commentary")
    build_clarity(prs)
    build_chart_slide(prs, "Performance Search (GSC)",
                       "Clics et impressions sur la période",
                       "chart_gsc_clicks_impressions", "gsc_commentary")
    build_table_slide(prs, "Top pages (GSC)",
                       "Pages d'atterrissage les plus performantes",
                       "table_top_pages")
    build_gmb_overview(prs)
    build_gmb_details(prs)
    build_backlinks_slide(prs, part=1)
    build_backlinks_slide(prs, part=2)
    build_final_summary_slide(prs)
    return prs


def generate_template(*, force: bool = False,
                      output: Path | None = None) -> Path:
    """Build or refresh the deck when ``force`` or version is stale."""
    target = output or resolve_template_path()
    if target.exists() and not force and not template_needs_rebuild(target):
        return target

    target.parent.mkdir(parents=True, exist_ok=True)
    prs = _build_presentation()
    prs.save(target)
    _write_template_version(target)
    return target


def main() -> None:
    args = _parse_args()
    target = resolve_template_path(args.output or None)
    force = bool(args.force)
    if args.force_if_stale and template_needs_rebuild(target):
        force = True

    if target.exists() and not force:
        print(
            f"Template already exists at {target}.\n"
            "Edit that file in PowerPoint to change layout or branding.\n"
            "Monthly reports use it as-is (placeholders only).\n"
            "To regenerate from code, run: python scripts/build_template.py --force",
            file=sys.stderr,
        )
        sys.exit(0)

    path = generate_template(force=True, output=target)
    print(f"Template generated at {path} (version {TEMPLATE_BUILD_VERSION})")


if __name__ == "__main__":
    main()
