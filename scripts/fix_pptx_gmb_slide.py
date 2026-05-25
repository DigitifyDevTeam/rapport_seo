"""One-off patch: replace GMB overview slide image + KPI values in a .pptx."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from src.reporting.pptx_report import (  # noqa: E402
    _PICTURE_MARGIN_RATIO,
    _fitted_picture_bounds,
)

_GMB_KPI_VALUES = {
    "702": "3 566",
    "79": "417",
    "218": "1 306",
    "405": "1 843",
}

# Crop box for GBP panel in the reference slide screenshot (full slide layout).
_REF_CROP = (0, 42, 478, None)  # height = full image height


def _crop_business_card(reference: Path, out: Path) -> Path:
    im = Image.open(reference)
    w, h = im.size
    left, top, right, bottom = _REF_CROP[0], _REF_CROP[1], _REF_CROP[2], _REF_CROP[3] or h
    cropped = im.crop((left, top, right, bottom))
    out.parent.mkdir(parents=True, exist_ok=True)
    cropped.save(out)
    return out


def _find_gmb_overview_slide(prs: Presentation):
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if any("Fiche publique" in t for t in texts) and any(
            "INTERACTIONS TOTALES" in t for t in texts
        ):
            return slide
    return None


def patch_pptx(pptx_path: Path, reference_image: Path) -> None:
    pptx_path = pptx_path.resolve()
    card_png = pptx_path.parent / f".{pptx_path.stem}_gmb_card_fix.png"
    _crop_business_card(reference_image, card_png)

    prs = Presentation(str(pptx_path))
    slide = _find_gmb_overview_slide(prs)
    if slide is None:
        raise SystemExit(f"No GMB overview slide in {pptx_path}")

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)  # noqa: SLF001
            fit_left, fit_top, fit_w, fit_h = _fitted_picture_bounds(
                left, top, width, height, card_png,
                margin_ratio=_PICTURE_MARGIN_RATIO,
            )
            slide.shapes.add_picture(str(card_png), fit_left, fit_top,
                                     width=fit_w, height=fit_h)
            break

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                key = run.text.strip()
                if key in _GMB_KPI_VALUES:
                    run.text = _GMB_KPI_VALUES[key]

    prs.save(str(pptx_path))
    print(f"Patched {pptx_path}")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path, nargs="+")
    parser.add_argument(
        "--reference",
        type=Path,
        required=True,
        help="Full slide screenshot with correct Fiche publique layout",
    )
    args = parser.parse_args()
    ref = args.reference.resolve()
    if not ref.is_file():
        raise SystemExit(f"Reference image not found: {ref}")
    for path in args.pptx:
        patch_pptx(path, ref)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
